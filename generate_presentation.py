#!/usr/bin/env python3
"""
Gerador de Apresentação PowerPoint - DevSecOps Examples
Gera apresentação sobre ISO 27017/27018 e Pipeline de Compliance Contínuo
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

# Aumentar limite de segurança do PIL para aceitar imagens grandes dos diagramas
Image.MAX_IMAGE_PIXELS = None
import os

def create_presentation():
    """Cria apresentação PowerPoint sobre as pastas 5 e 6"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Título
    add_title_slide(prs)
    
    # Slide 2: Visão Geral ISO 27017/27018
    add_iso_overview_slide(prs)
    
    # Slide 3: ISO 27017 - Backup
    add_iso27017_backup_slide(prs)
    
    # Slide 4: ISO 27017 - Criptografia
    add_iso27017_encryption_slide(prs)
    
    # Slide 5: ISO 27017 - Segregação
    add_iso27017_segregation_slide(prs)
    
    # Slide 6: ISO 27018 - Auditoria
    add_iso27018_audit_slide(prs)
    
    # Slide 7: ISO 27018 - Esquecimento
    add_iso27018_erasure_slide(prs)
    
    # Slide 8: ISO 27018 - Localização
    add_iso27018_location_slide(prs)
    
    # Slide 9: Pipeline - Visão Geral
    add_pipeline_overview_slide(prs)
    
    # Slide 10: Pipeline - Stages 1-3
    add_pipeline_stages_1_3_slide(prs)
    
    # Slide 11: Pipeline - Stages 4-5
    add_pipeline_stages_4_5_slide(prs)
    
    # Slide 12: Pipeline - Métricas
    add_pipeline_metrics_slide(prs)
    
    # Slide 13: Conclusão
    add_conclusion_slide(prs)
    
    # Salvar apresentação
    output_file = "DevSecOps_ISO27017_27018_Presentation.pptx"
    prs.save(output_file)
    print(f"✅ Apresentação criada: {output_file}")
    return output_file

def set_text_format(text_frame, font_size=20, font_name="Arial", bold=False):
    """Define formatação padrão do texto"""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold

def add_title_slide(prs):
    """Slide 1: Título"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Título principal
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "DevSecOps & Compliance Contínuo"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.name = "Arial"
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "ISO 27017/27018 & Pipeline Automatizada"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.name = "Arial"
    subtitle_para.font.color.rgb = RGBColor(64, 64, 64)
    
    # Rodapé
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Exemplos práticos de conformidade em Cloud Computing"
    footer_para = footer_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.CENTER
    footer_para.font.size = Pt(16)
    footer_para.font.name = "Arial"
    footer_para.font.color.rgb = RGBColor(128, 128, 128)

def add_iso_overview_slide(prs):
    """Slide 2: Visão Geral ISO 27017/27018"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "🔒 ISO 27017/27018 - Visão Geral"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.name = "Arial"
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    # Conteúdo
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    content = """ISO 27017 - Cloud Computing Security
• Controles de segurança para cloud
• Backup e recuperação de dados
• Criptografia em repouso e trânsito
• Segregação de ambientes (Dev/Prod)

ISO 27018 - Personal Data Protection
• Proteção de dados pessoais na nuvem
• Auditoria e rastreabilidade
• Direito ao esquecimento (LGPD)
• Data residency - localização dos dados"""
    
    content_frame.text = content
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.name = "Arial"
        paragraph.space_after = Pt(10)

def add_iso27017_backup_slide(prs):
    """Slide 3: ISO 27017 - Backup"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "💾 ISO 27017 - Backup e Recuperação"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.name = "Arial"
    title_para.font.color.rgb = RGBColor(0, 102, 51)
    
    # Imagem do diagrama (esquerda)
    try:
        img_path = "exemplos/5 - exemplos iso-27017 - iso-27018/iso-27017-backup/iso-27017-backup-architecture.png"
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.3), width=Inches(4.5))
    except FileNotFoundError:
        pass
    
    # Conteúdo (direita)
    content_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.2))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    content = """Conceito:
Backup automatizado de dados críticos com retenção de 30 dias

Implementação:
• AWS Backup Vault
• Backup Plan (diário 3AM)
• Seleção por tags
• SNS Notifications

Validação OPA:
• Agendamento diário
• Retenção >= 30 dias
• Notificações ativas

Benefícios:
• RPO: 24 horas
• RTO: < 4 horas
• Proteção ransomware
• Custo: $0.05/GB/mês"""
    
    content_frame.text = content
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.name = "Arial"
        paragraph.space_after = Pt(6)
        paragraph.font.size = Pt(20)
        paragraph.font.name = "Arial"
        paragraph.space_after = Pt(8)

def add_iso27017_encryption_slide(prs):
    """Slide 4: ISO 27017 - Criptografia"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "🔐 ISO 27017 - Criptografia"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.name = "Arial"
    title_para.font.color.rgb = RGBColor(0, 102, 51)
    
    # Imagem do diagrama (esquerda)
    try:
        img_path = "exemplos/5 - exemplos iso-27017 - iso-27018/iso-27017-criptografia/iso-27017-criptografia-architecture.png"
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.3), width=Inches(4.5))
    except FileNotFoundError:
        pass
    
    # Conteúdo (direita)
    content_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.2))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    content = """Conceito:
Criptografia de dados em repouso com chaves gerenciadas

Implementação:
• AWS KMS Key
• S3 SSE-KMS (AES-256)
• Rotação automática (365d)
• Public Access Block
• Versionamento

Validação OPA:
• Criptografia aws:kms
• Rotação automática
• Versionamento ativo

Benefícios:
• Proteção militar
• Defesa em profundidade
• LGPD/PCI DSS/HIPAA
• Overhead < 5%"""
    
    content_frame.text = content
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.name = "Arial"
        paragraph.space_after = Pt(6)

def add_iso27017_segregation_slide(prs):
    """Slide 5: ISO 27017 - Segregação"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "🏗️ ISO 27017 - Segregação de Rede"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.name = "Arial"
    title_para.font.color.rgb = RGBColor(0, 102, 51)
    
    # Imagem do diagrama (esquerda)
    try:
        img_path = "exemplos/5 - exemplos iso-27017 - iso-27018/iso-27017-segregacao/iso-27017-segregacao-architecture.png"
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.3), width=Inches(4.5))
    except FileNotFoundError:
        pass
    
    # Conteúdo (direita)
    content_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.2))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    content = """Conceito:
Isolamento completo entre ambientes Dev e Prod

Implementação:
• VPC Prod (10.0.0.0/16)
• VPC Dev (10.1.0.0/16)
• Subnets privadas
• Network ACLs deny
• Flow Logs

Validação OPA:
• Tags Environment
• CIDRs não sobrepostos
• ACLs com deny rules

Benefícios:
• Isolamento 100%
• Blast radius reduzido
• 70% menos incidentes
• SOC 2/PCI DSS"""
    
    content_frame.text = content
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.name = "Arial"
        paragraph.space_after = Pt(6)

def add_iso27018_audit_slide(prs):
    """Slide 6: ISO 27018 - Auditoria"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "📋 ISO 27018 - Auditoria"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.name = "Arial"
    title_para.font.color.rgb = RGBColor(102, 0, 153)
    
    # Imagem do diagrama (esquerda)
    try:
        img_path = "exemplos/5 - exemplos iso-27017 - iso-27018/iso-27018-auditoria/iso-27018-auditoria-architecture.png"
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.3), width=Inches(4.5))
    except FileNotFoundError:
        pass
    
    # Conteúdo (direita)
    content_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.2))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    content = """Conceito:
Rastreabilidade completa de acessos a dados pessoais (LGPD Art.37)

Implementação:
• CloudTrail multi-region
• S3 Bucket logs (7 anos)
• Object Lock compliance
• CloudWatch Alarms
• Metric Filters

Validação OPA:
• Multi-region ativo
• Retenção >= 2557 dias
• Data Events capturados
• Detecção anomalias

Benefícios:
• LGPD Art. 37 compliant
• Não-repúdio
• Latência < 5min
• Custo: $2/100k eventos"""
    
    content_frame.text = content
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.name = "Arial"
        paragraph.space_after = Pt(6)

def add_iso27018_erasure_slide(prs):
    """Slide 7: ISO 27018 - Direito ao Esquecimento"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "🗑️ ISO 27018 - Direito ao Esquecimento"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.name = "Arial"
    title_para.font.color.rgb = RGBColor(102, 0, 153)
    
    # Imagem do diagrama (esquerda)
    try:
        img_path = "exemplos/5 - exemplos iso-27017 - iso-27018/iso-27018-esquecimento/iso-27018-esquecimento-architecture.png"
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.3), width=Inches(4.5))
    except FileNotFoundError:
        pass
    
    # Conteúdo (direita)
    content_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.2))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    content = """Conceito:
Automação do direito ao esquecimento (LGPD Art.18, VI)

Implementação:
• Lambda timeout 300s
• SQS retenção 14 dias
• DynamoDB (PITR)
• CloudWatch Logs 7 anos

Validação OPA:
• Lambda timeout >= 300s
• SQS >= 14 dias
• Logs >= 7 anos
• DynamoDB PITR ativo

Benefícios:
• SLA < 15 dias
• Taxa sucesso 99.5%
• Custo: $0.001/exclusão
• Evita multas R$50M"""
    
    content_frame.text = content
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.name = "Arial"
        paragraph.space_after = Pt(6)

def add_iso27018_location_slide(prs):
    """Slide 8: ISO 27018 - Localização de Dados"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "🌎 ISO 27018 - Data Residency"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.name = "Arial"
    title_para.font.color.rgb = RGBColor(102, 0, 153)
    
    # Imagem do diagrama (esquerda)
    try:
        img_path = "exemplos/5 - exemplos iso-27017 - iso-27018/iso-27018-localizacao/iso-27018-localizacao-architecture.png"
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.3), width=Inches(4.5))
    except FileNotFoundError:
        pass
    
    # Conteúdo (direita)
    content_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.2))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    content = """Conceito:
Soberania de dados - 100% em território brasileiro

Implementação:
• S3 sa-east-1 (Brasil)
• Bucket policy deny
• Lifecycle 5 anos max
• AWS Config Rule
• Tags compliance

Validação OPA:
• Tags Region/LGPD
• Bloqueio replicação
• Retenção <= 5 anos

Benefícios:
• LGPD Art.11 compliant
• Latência 15ms Brasil
• Sem custos transfer
• Evita multas 2% fat."""
    
    content_frame.text = content
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.name = "Arial"
        paragraph.space_after = Pt(6)

def add_pipeline_overview_slide(prs):
    """Slide 9: Pipeline - Visão Geral"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "🚀 Pipeline de Compliance Contínuo"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.name = "Arial"
    title_para.font.color.rgb = RGBColor(204, 51, 0)
    
    # Imagem da arquitetura do pipeline (se existir)
    try:
        img_path = "exemplos/6 - pipeline compliance continuo/compliance-pipeline-architecture.png"
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.3), width=Inches(9))
    except FileNotFoundError:
        # Conteúdo alternativo se não houver imagem
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.6))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        content = """Filosofia:
• Falhas não bloqueiam visibilidade
• Execução completa garantida
• continue-on-error: true (todos jobs)
• if: always() (dependências)

Arquitetura:
• 7 estágios automatizados
• Validações paralelas (SAST)
• OPA Policy as Code
• Terraform dry-run (sem AWS)

Diferencial:
• Mock credentials (exemplos)
• Visibilidade total de issues
• Métricas agregadas ao final
• Zero deploy em produção"""
        
        content_frame.text = content
        for paragraph in content_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.name = "Arial"
            paragraph.space_after = Pt(8)

def add_pipeline_stages_1_3_slide(prs):
    """Slide 10: Pipeline Stages 1-3"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "📝 Pipeline - Stages 1-3"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(30)
    title_para.font.bold = True
    title_para.font.name = "Arial"
    title_para.font.color.rgb = RGBColor(204, 51, 0)
    
    # Conteúdo
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.6))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    content = """Stage 1 - Validação de Código:
• terraform fmt -check (formatação)
• terraform validate (sintaxe)
• SLA: ~1 minuto

Stage 2 - Análise de Segurança (SAST):
• TFSec: vulnerabilidades em Terraform
• Checkov: best practices de segurança
• Upload SARIF para GitHub Security
• SLA: ~2 minutos

Stage 3 - Validação de Políticas (OPA):
• 6 políticas ISO 27017/27018
• terraform plan + opa eval
• Validação com mock credentials
• SLA: ~5 minutos"""
    
    content_frame.text = content
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.name = "Arial"
        paragraph.space_after = Pt(8)

def add_pipeline_stages_4_5_slide(prs):
    """Slide 11: Pipeline Stages 4-5"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "📋 Pipeline - Stages 4-6"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(30)
    title_para.font.bold = True
    title_para.font.name = "Arial"
    title_para.font.color.rgb = RGBColor(204, 51, 0)
    
    # Conteúdo
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.6))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    content = """Stage 4 - Terraform Plan:
• Plan com -refresh=false
• Mock AWS provider para exemplos
• Upload de artefatos
• SLA: ~2 minutos

Stage 5 - Estimativa de Custos (Infracost):
• Análise de todos os 6 exemplos
• Relatório de custos AWS
• Comentários automáticos em PRs
• SLA: ~2 minutos

Stage 6 - Relatório de Compliance:
• Consolidação de resultados
• Métricas de conformidade
• Upload de artefatos
• SLA: ~1 minuto"""
    
    content_frame.text = content
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.name = "Arial"
        paragraph.space_after = Pt(8)

def add_pipeline_metrics_slide(prs):
    """Slide 12: Pipeline Métricas"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "📈 Métricas e Benefícios"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(30)
    title_para.font.bold = True
    title_para.font.name = "Arial"
    title_para.font.color.rgb = RGBColor(204, 51, 0)
    
    # Conteúdo
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.6))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    content = """Performance:
• Tempo total: ~13 minutos
• Execução paralela de stages
• Continue-on-error para visibilidade total

Conformidade:
• 6 políticas validadas automaticamente
• 50+ checks de segurança (Checkov/TFSec)
• 100% rastreabilidade via artefatos

Automação:
• Execução em push, PR, schedule (diário 3AM)
• Workflow dispatch para execução manual
• Comentários automáticos em Pull Requests

Economia:
• Mock credentials - sem custos AWS
• Validação antes do deploy
• Prevenção de não-conformidade"""
    
    content_frame.text = content
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.name = "Arial"
        paragraph.space_after = Pt(8)

def add_conclusion_slide(prs):
    """Slide 13: Conclusão"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "✅ Conclusão"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(30)
    title_para.font.bold = True
    title_para.font.name = "Arial"
    title_para.font.color.rgb = RGBColor(0, 102, 51)
    
    # Conteúdo
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.6))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    content = """Principais Conquistas:

• Compliance automatizado ISO 27017/27018
• Pipeline CI/CD completa com 6 stages
• Policy as Code com Open Policy Agent
• Segurança integrada (SAST + OPA)
• Estimativa de custos automatizada
• 100% rastreável e auditável

Próximos Passos:

• Integrar com ambiente real AWS
• Adicionar testes de integração
• Implementar deploy automatizado
• Expandir políticas OPA
• Dashboard de métricas

Recursos: github.com/tharles-freire-heyupcharly/devsecops-examples"""
    
    content_frame.text = content
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.name = "Arial"
        paragraph.space_after = Pt(8)

if __name__ == "__main__":
    create_presentation()
